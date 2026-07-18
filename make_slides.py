#!/usr/bin/env python3
"""Generate the CryoEMAgent AAAI reference deck.
Style: white background, black text, Times New Roman. 4 core slides + title + architecture figure.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)
FONT = "Times New Roman"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def white_bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = WHITE


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def style_run(r, size, bold=False, italic=False, color=BLACK):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def set_para(p, text, size, bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_after=6, level=0):
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    style_run(r, size, bold, italic, color)
    return p


def title_block(slide, kicker, title):
    tb, tf = textbox(slide, 0.7, 0.45, 11.9, 1.1)
    p = tf.paragraphs[0]
    set_para(p, kicker, 14, bold=True, color=GREY, space_after=2)
    p2 = tf.add_paragraph()
    set_para(p2, title, 30, bold=True, space_after=0)
    # underline rule
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.62), Inches(11.9), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLACK; ln.line.fill.background()


def bullets(slide, items, top=1.95, left=0.8, width=11.7, height=5.2, base=19):
    tb, tf = textbox(slide, left, top, width, height)
    first = True
    for it in items:
        lvl = it[0]; txt = it[1]
        bold = it[2] if len(it) > 2 else False
        italic = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        glyph = "•   " if lvl == 0 else "–   "
        size = base if lvl == 0 else base - 2
        set_para(p, glyph + txt, size, bold=bold, italic=italic,
                 space_after=7 if lvl == 0 else 4, level=lvl)
    return tb


def box(slide, l, t, w, h, lines, fill=WHITE, fsize=13, bold_first=True):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = BLACK; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_para(p, ln, fsize if i else fsize + 1, bold=(bold_first and i == 0),
                 align=PP_ALIGN.CENTER, space_after=1)
    return sp


def down_arrow(slide, cx, top, h, label=None):
    w = 0.5
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx - w / 2), Inches(top), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = BLACK; a.line.fill.background(); a.shadow.inherit = False
    if label:
        tb, tf = textbox(slide, cx + 0.4, top + h / 2 - 0.22, 5.2, 0.5, MSO_ANCHOR.MIDDLE)
        set_para(tf.paragraphs[0], label, 12, italic=True, color=GREY)


def comparison_table(slide, top=2.0, left=0.7, width=11.9):
    rows = [
        ("Aspect", "CryoWizard / Cryo-IEF", "CryoEMAgent (ours)"),
        ("Particle selection", "Fixed threshold (set once)", "ACPS — adaptive, closed-loop (resolution feedback)"),
        ("Safety", "None", "Keep-best guarantee (never worse than baseline)"),
        ("Failure recovery", "Fixed / same-recipe retries", "Domain-Grounded Reflection (causal failure graph)"),
        ("Operation", "On the GPU server", "Laptop-driven over MCP-over-SSH"),
        ("Decision model", "None (hard-coded recipe)", "SPA formalized as an MDP"),
        ("Evaluation", "Resolution only", "+ CryoPipelineQA benchmark + VLM grounding audit"),
    ]
    nr, nc = len(rows), 3
    gt = slide.shapes.add_table(nr, nc, Inches(left), Inches(top),
                                Inches(width), Inches(0.55 * nr)).table
    # plain black grid, no fills, no banding ("No Style, Table Grid")
    tblPr = gt._tbl.tblPr
    for el in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(el)
    sid = tblPr.makeelement(qn('a:tableStyleId'), {})
    sid.text = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    tblPr.append(sid)
    gt.first_row = True
    gt.horz_banding = False
    gt.columns[0].width = Inches(2.5)
    gt.columns[1].width = Inches(3.9)
    gt.columns[2].width = Inches(width - 6.4)
    for r in range(nr):
        for c in range(nc):
            cell = gt.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = rows[r][c]
            style_run(run, 13, bold=(r == 0 or c == 0))
    return gt


# ---------------------------------------------------------------- TITLE
s = prs.slides.add_slide(BLANK); white_bg(s)
tb, tf = textbox(s, 1.0, 2.3, 11.3, 2.6, MSO_ANCHOR.TOP)
set_para(tf.paragraphs[0], "CryoEMAgent", 44, bold=True, space_after=8)
set_para(tf.add_paragraph(),
         "An Autonomous LLM Operator for Cryo-EM Single-Particle Analysis",
         24, space_after=18)
set_para(tf.add_paragraph(),
         "Adaptive, principled, and evaluable decision-making — beyond fixed-threshold automation",
         16, italic=True, color=GREY, space_after=24)
set_para(tf.add_paragraph(), "Yaswanth Devavarapu   ·   Rakshitha Ireddi   ·   Xu Lab, Carnegie Mellon University",
         15, space_after=2)
set_para(tf.add_paragraph(), "Mentor: Mei Yuan", 14, color=GREY)

# ---------------------------------------------------------------- 1. RESEARCH GAP
s = prs.slides.add_slide(BLANK); white_bg(s)
title_block(s, "MOTIVATION", "Research Gap")
bullets(s, [
    (0, "Cryo-EM SPA is decision-heavy: curation, picking, 2D/class selection, and refinement each need expert judgement — the human is the throughput bottleneck.", True),
    (0, "Existing automation handles the pipeline, not the decision process:"),
    (1, "CryoWizard / Cryo-IEF — particle selection by FIXED thresholds."),
    (1, "Structura (2025) — fixed recipe; retries the SAME recipe on failure."),
    (1, "None formalize the decisions, adapt selection from outcomes, or benchmark the reasoning."),
    (0, "Open questions no prior system addresses:"),
    (1, "Can SPA decision-making be posed as a principled, optimisable framework?"),
    (1, "Can particle selection ADAPT in a closed loop instead of using a fixed cutoff?"),
    (1, "Do vision-language critics actually SEE the images, or only read the metrics?"),
    (0, "Gap: treat SPA decisions as an adaptive, evaluable process — not hard-coded thresholds.", True),
])

# ---------------------------------------------------------------- 2. METHOD
s = prs.slides.add_slide(BLANK); white_bg(s)
title_block(s, "APPROACH", "Method — Five New Contributions")
bullets(s, [
    (0, "CryoEMAgent: an LLM-driven autonomous CryoSPARC operator, run from a laptop over MCP-over-SSH, with zero human checkpoints.", True),
    (0, "(1) SPA-as-MDP — first formalization of single-particle analysis as a Markov Decision Process (states, actions, transitions, reward)."),
    (0, "(2) ACPS — closed-loop ADAPTIVE particle selection:  threshold ← threshold − α·(resolution − target), with a keep-best safety guarantee (never worse than baseline). Replaces fixed thresholds."),
    (0, "(3) Domain-Grounded Reflection — failure recovery over a cryo-EM CAUSAL FAILURE GRAPH (not generic Reflexion)."),
    (0, "(4) CryoPipelineQA — the first BENCHMARK for evaluating LLM reasoning on cryo-EM processing decisions."),
    (0, "(5) VLM visual-grounding audit — does the model see the micrograph / 2D images, or only read the text metrics?"),
    (0, "Architecture overview in Figure 1 (next slide).", False, True),
])

# ---------------------------------------------------------------- FIGURE: ARCHITECTURE
s = prs.slides.add_slide(BLANK); white_bg(s)
title_block(s, "FIGURE 1", "System Architecture")
cx = 13.333 / 2
bw = 8.4; bl = cx - bw / 2
box(s, bl, 1.85, bw, 1.15, [
    "Laptop  (Windows, no GPU)  —  CryoEMAgent",
    "ReAct Planner   ·   ACPS (adaptive selection)   ·   Quality Critic   ·   Reasoning / Audit log",
], fsize=12.5)
down_arrow(s, cx, 3.05, 0.7, "MCP-over-SSH   (JSONL stdio, two-hop jump host)")
box(s, bl, 3.8, bw, 1.0, [
    "GPU Server",
    "Orchestrator (W1 blob arm / W2 template arm)   ·   MCP dispatcher",
], fsize=12.5)
down_arrow(s, cx, 4.85, 0.6, "REST API  (localhost:39000)")
box(s, bl, 5.5, bw, 1.15, [
    "CryoSPARC v4.7.1",
    "import → motion → CTF → auto-curate → pick → extract → 2D → select → ab-initio → refine",
], fsize=12)
# feedback loop arrow (right side, bottom -> ACPS)
fb = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(bl + bw + 0.15), Inches(2.0), Inches(0.45), Inches(4.55))
fb.fill.solid(); fb.fill.fore_color.rgb = GREY; fb.line.fill.background(); fb.shadow.inherit = False
tbf, tff = textbox(s, bl + bw + 0.65, 4.0, 2.4, 1.0, MSO_ANCHOR.MIDDLE)
set_para(tff.paragraphs[0], "resolution feedback", 11.5, italic=True, color=GREY, space_after=1)
set_para(tff.add_paragraph(), "(ACPS closed loop)", 11.5, italic=True, color=GREY)

# ---------------------------------------------------------------- 3. COMPARISON TABLE
s = prs.slides.add_slide(BLANK); white_bg(s)
title_block(s, "POSITIONING", "CryoEMAgent vs. Prior Automation")
comparison_table(s, top=2.05)
tb, tf = textbox(s, 0.72, 6.7, 11.9, 0.6)
set_para(tf.paragraphs[0],
         "Note: the 3.21 Å result used CryoSPARC's built-in (fixed-style) selection; ACPS is our adaptive replacement — the head-to-head ACPS-vs-fixed result is in progress.",
         12, italic=True, color=GREY)

# ---------------------------------------------------------------- 4. EXPERIMENTS
s = prs.slides.add_slide(BLANK); white_bg(s)
title_block(s, "RESULTS", "Experiments")
bullets(s, [
    (0, "“→” means “improves / scales to” — a transition, NOT a method comparison.", False, True),
    (0, "Autonomous operator — feasibility (uses CryoSPARC's built-in selection, NOT ACPS):", True),
    (1, "3.21 Å on full EMPIAR-10288 (CB1-GPCR, 2,756 movies, 477 GB), zero human clicks, ~20.7 h."),
    (1, "Auto-curation rejected 21%; 620,174 particles from ~2M picks; auto-recovered from a GPU out-of-memory event."),
    (1, "Scaling check, pilot → full (same method, more data): 5,358 → 620,174 particles; 9.12 → 3.21 Å."),
    (0, "ACPS (contribution 2) — built + validated:", True),
    (1, "Control law + keep-best safety; selection knob monotonic (looser keeps more: top-30% → 938k, top-70% → 1.40M particles); resolution feedback validated; 16 unit tests; end-to-end run on EMPIAR-10146."),
    (1, "ACPS baseline = fixed-threshold selection (CryoWizard / Cryo-IEF style). Head-to-head ACPS vs fixed: in progress [NTD]."),
    (0, "In progress: ACPS winning run (EMPIAR-10028 ribosome, C1) · CryoPipelineQA · Domain-Grounded Reflection · VLM grounding audit · MDP write-up.", True),
], base=15)

# ---------------------------------------------------------------- 4. CONCLUSION
s = prs.slides.add_slide(BLANK); white_bg(s)
title_block(s, "FINDINGS", "Conclusion & Findings")
bullets(s, [
    (0, "Demonstrated: cryo-EM SPA can run fully autonomously from a laptop to a 3.21 Å near-atomic GPCR map — feasibility established.", True),
    (0, "Contributed (built): ACPS, a closed-loop adaptive particle-selection algorithm with a provable keep-best safety property — a principled replacement for fixed thresholds.", True),
    (0, "Findings so far: autonomous remote SPA is viable; adaptive selection with a safety guarantee is implementable and stable."),
    (0, "Honest limitations:"),
    (1, "The 3.21 Å result validates the autonomous OPERATOR; it uses CryoSPARC's selection, not ACPS."),
    (1, "ACPS still needs its winning result; four of five contributions are designed but not yet executed; benchmarking is single-protein so far."),
    (0, "Next 2–4 weeks: ACPS result on a C1 dataset → MDP formalization → CryoPipelineQA-lite.", True),
])

out = "CryoEMAgent_AAAI_slides.pptx"
prs.save(out)
print("WROTE", out, "slides:", len(prs.slides._sldIdLst))
